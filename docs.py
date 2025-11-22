import docx
from pptx import Presentation
import random
import os
import subprocess
from ai import ai_reply
from tts import speak

# Ensure documents folder exists
SAVE_DIR = os.path.join(os.getcwd(), "documents")
os.makedirs(SAVE_DIR, exist_ok=True)

def open_file(filepath):
    """Open file with the default application based on OS"""
    try:
        if os.name == "nt":  # Windows
            os.startfile(filepath)
        elif os.name == "posix":  # macOS / Linux
            subprocess.call(("open" if sys.platform == "darwin" else "xdg-open", filepath))
    except Exception as e:
        print(f"⚠️ Could not open file: {e}")

def create_word_doc(topic, log_func=None):
    msg_start = f"📄 Creating Word document on {topic}..."
    speak(msg_start, log_func)
        
    ai_text = ai_reply(f"Write a detailed report on {topic}.")
    doc = docx.Document()
    doc.add_heading(f"Report on {topic}", 0)
    for line in ai_text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)

    filename = f"Jarvis_Document_{random.randint(1,1000)}.docx"
    filepath = os.path.join(SAVE_DIR, filename)
    doc.save(filepath)

    msg = f"✅ Word document created: {filepath}"
    speak(msg, log_func)

    open_file(filepath)  # 🔹 auto-open


def create_ppt(topic, log_func=None):
    msg_start = f"📊 Creating PowerPoint on {topic}..."
    speak(msg_start, log_func)

    ai_text = ai_reply(f"Create a 5-slide presentation on {topic}.")
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slides = ai_text.split("Slide")
    for slide in slides:
        if slide.strip():
            lines = [line.strip() for line in slide.split("\n") if line.strip()]
            slide_obj = prs.slides.add_slide(slide_layout)
            slide_obj.shapes.title.text = lines[0]
            tf = slide_obj.placeholders[1].text_frame
            for line in lines[1:]:
                tf.add_paragraph().text = line

    filename = f"Jarvis_Presentation_{random.randint(1,1000)}.pptx"
    filepath = os.path.join(SAVE_DIR, filename)
    prs.save(filepath)
    
    msg = f"✅ PowerPoint created: {filepath}"
    speak(msg, log_func)
 
    open_file(filepath)  # 🔹 auto-open
